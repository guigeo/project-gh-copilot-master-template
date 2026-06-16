"""Builder de apresentação: aplica specs de conteúdo num template via layouts.

Parte de um template `.pptx` de marca (ou do template padrão do python-pptx,
se nenhum for informado) e adiciona slides pelos layouts existentes,
preenchendo placeholders — sem posicionar shapes no olho.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

from .content import SlideBullets

# Índices dos layouts no template padrão do python-pptx.
LAYOUT_CAPA = 0  # Slide de título
LAYOUT_TITULO_CORPO = 1  # Título e conteúdo
LAYOUT_SO_TITULO = 5  # Apenas título (bom para gráfico/tabela)


def nova_apresentacao(template: str | Path | None = None) -> PresentationType:
    """Abre o template de marca; sem template, usa o padrão do python-pptx."""
    return Presentation(str(template)) if template else Presentation()


def add_capa(prs: PresentationType, titulo: str, subtitulo: str = ""):
    """Adiciona o slide de capa preenchendo título e subtítulo (placeholders)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CAPA])
    slide.shapes.title.text = titulo
    if subtitulo and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitulo
    return slide


def add_bullets(prs: PresentationType, spec: SlideBullets):
    """Adiciona um slide de tópicos a partir de uma spec pura."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITULO_CORPO])
    slide.shapes.title.text = spec.titulo

    corpo = slide.placeholders[1].text_frame
    corpo.clear()
    for i, texto in enumerate(spec.bullets):
        paragrafo = corpo.paragraphs[0] if i == 0 else corpo.add_paragraph()
        paragrafo.text = texto
    return slide


def add_slide_visual(prs: PresentationType, titulo: str):
    """Adiciona um slide só com título, pronto para receber gráfico/tabela nativo."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SO_TITULO])
    slide.shapes.title.text = titulo
    return slide


def salvar(prs: PresentationType, caminho: str | Path) -> None:
    """Persiste o `.pptx` (I/O na borda); cria o diretório-pai se faltar."""
    destino = Path(caminho)
    destino.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destino))
