"""Gráficos e tabelas NATIVOS a partir de DataFrame (editáveis no PowerPoint).

Alimente sempre com um DataFrame já agregado pelo pandas; o slide mostra
resultado, não dado bruto.
"""

from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

import pandas as pd


def add_grafico_barras(
    slide,
    df: pd.DataFrame,
    categoria: str,
    valor: str,
    titulo: str | None = None,
):
    """Adiciona um gráfico de colunas nativo (categoria x valor) ao slide."""
    dados = CategoryChartData()
    dados.categories = df[categoria].astype(str).tolist()
    dados.add_series(valor, df[valor].tolist())

    grafico_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1.5),
        Inches(8),
        Inches(5),
        dados,
    )
    chart = grafico_frame.chart
    if titulo:
        chart.has_title = True
        chart.chart_title.text_frame.text = titulo
    return grafico_frame


def add_tabela(slide, df: pd.DataFrame, topo: float = 1.5):
    """Adiciona uma tabela nativa com o conteúdo do DataFrame (cabeçalho na linha 0).

    Formate números no DataFrame antes de chamar — a tabela exibe texto.
    """
    n_linhas, n_colunas = df.shape[0] + 1, df.shape[1]
    tabela_frame = slide.shapes.add_table(
        n_linhas,
        n_colunas,
        Inches(0.5),
        Inches(topo),
        Inches(9),
        Inches(0.4 * n_linhas),
    )
    tabela = tabela_frame.table

    for j, coluna in enumerate(df.columns):
        tabela.cell(0, j).text = str(coluna)
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            tabela.cell(i + 1, j).text = str(df.iat[i, j])
    return tabela_frame
